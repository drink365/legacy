# modules/strategy_pack.py
# -*- coding: utf-8 -*-
from io import BytesIO
from pathlib import Path
import zipfile

# ---------- 預設品牌參數 ----------
BRAND = {
    "blue": {"bg": "1A365D", "gold": "D4AF37"},
    "red":  {"bg": "A40E26", "gold": "E8C173"},
}

DEFAULTS = {
    "author": "Grace Huang｜永傳家族辦公室創辦人",
    "footer_text": "《影響力》傳承策略平台｜永傳家族辦公室  |  https://gracefo.com  |  聯絡信箱：123@gracefo.com",
    "closing_quote": "讓智慧，承載愛的延續。\n永傳家族辦公室，與您一起，擘劃幸福的未來。",
    "subtitle_blue": "專家洞見 × 智能科技 × 幸福傳承",
    "subtitle_red": "讓智慧，承載愛的延續。",
}

SLIDES_OUTLINE = [
    ("封面", "永傳家族辦公室｜商業策略藍圖", ""),  # 副標視版本帶入
    ("品牌願景與定位", "願景：AI × 傳承 × 溫度", "以專業與科技，守護家族的幸福延續"),
    ("市場洞察與痛點", "痛點：複雜法稅、跨境、爭產風險", "機會：以AI與顧問服務簡化決策"),
    ("目標客群與核心需求", "B2C：高資產家庭；B2B：顧問夥伴", "需求：合法節稅、結構清楚、安心交棒"),
    ("使用者旅程（摘要）", "探索→理解→信任→決策→推薦", "AI體驗 + 顧問解讀 + 提案落地"),
    ("品牌價值鏈", "教育→體驗→轉化", "公開講座/AI試算 → 案例解讀 → 方案落地"),
    ("三階段成長策略 - 總覽", "Phase1 啟動｜Phase2 成長｜Phase3 生態", "短中長期路線圖"),
    ("Phase 1：啟動期（0–6個月）", "打造 AI 傳承入口、導流與預約", "Demo工具、名單收集、Line OA流程"),
    ("Phase 2：成長期（6–18個月）", "顧問賦能與商業化", "顧問登入、白標工具、內容訂閱"),
    ("Phase 3：生態期（18個月+）", "平台擴張與國際化", "API合作、家族帳戶、雙語市場"),
    ("營收模型", "B2C 顧問費＋保單佣金", "B2B 工具授權＋課程訂閱＋白標"),
    ("關鍵指標", "名單、轉化率、留存、MRR", "分階段設定KPI與里程碑"),
    ("品牌願景宣言", "不只延續資產，更延續愛與價值", "溫度 × 專業 × 永續"),
    ("行動計畫與里程碑", "接下來90天重點", "工具上線、案例包、導入顧問"),
    ("封底", "Thank You", ""),  # 引言用 closing_quote
]

REPORT_PARAGRAPHS = [
    ("Executive Summary", "永傳家族辦公室以 AI × 傳承 × 溫度 為核心，打造高資產家庭與專業顧問的雙向平台。透過教育→體驗→轉化的價值鏈，實現短期變現、中期賦能、長期生態的成長藍圖。"),
    ("市場洞察與機會", "面對複雜法稅、跨境資產與家族治理課題，市場需要兼具合法合規與情感理解的整合型解方。AI 工具讓認知與決策更高效，顧問服務讓信任更可被感受。"),
    ("目標客群與需求", "B2C 客群（企業主、一代創辦人、高資產家庭）重視安全、合規與安心交棒；B2B 客群（保險與財稅顧問）需要提升專業形象、教育工具與成交效率。"),
    ("使用者旅程（UX）", "探索→理解→信任→決策→推薦；以 AI 試算引發覺察，透過顧問解讀會建立信任，進而完成保單/信託/贈與等結構安排。"),
    ("品牌價值鏈", "公開講座/社群內容/AI工具 → 案例解讀/提案PDF → 方案落地/長期陪伴；讓體驗自然銜接轉化。"),
    ("三階段策略藍圖", "Phase1：啟動（AI入口與導流）｜Phase2：成長（顧問賦能與商業化）｜Phase3：生態（API合作與國際化）。"),
    ("營收模型", "B2C：顧問費＋保單佣金；B2B：工具授權＋課程訂閱；B2B2C：企業/白標合作。"),
    ("願景與收束", "讓智慧承載愛的延續，讓傳承成為幸福的流動。"),
]

def _hex_to_rgb(hexstr: str):
    hexstr = hexstr.strip().lstrip('#')
    return tuple(int(hexstr[i:i+2], 16) for i in (0,2,4))

def build_pptx_bytes(theme: str, *, author: str, footer_text: str, closing_quote: str,
                     logo_bytes: bytes | None, subtitle: str, with_tagline_bg: bool) -> bytes:
    # lazy import
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    pres = Presentation()
    pres.slide_width, pres.slide_height = Inches(13.33), Inches(7.5)
    bg = _hex_to_rgb(BRAND[theme]["bg"])
    gold = _hex_to_rgb(BRAND[theme]["gold"])
    TITLE_FONT = "Noto Serif TC"
    BODY_FONT = "Noto Sans TC"

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg)

    def footer(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*gold); shape.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.4))
        p = tb.text_frame.paragraphs[0]; p.text = footer_text; p.font.size = Pt(12); p.font.name = BODY_FONT

    for idx, (title, line1, line2) in enumerate(SLIDES_OUTLINE):
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        set_bg(slide)

        if idx == 0:
            # logo
            if logo_bytes:
                from PIL import Image
                import io
                img = Image.open(BytesIO(logo_bytes))
                bio = io.BytesIO(); img.save(bio, format="PNG"); bio.seek(0)
                slide.shapes.add_picture(bio, Inches(0.5), Inches(0.5), height=Inches(0.9))
            # title
            tbox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.2))
            p = tbox.text_frame.paragraphs[0]; p.text = line1; p.font.size = Pt(40); p.font.bold = True
            p.font.name = TITLE_FONT; p.font.color.rgb = RGBColor(*gold)
            # subtitle
            sbox = slide.shapes.add_textbox(Inches(1.2), Inches(3.3), Inches(11.0), Inches(0.8))
            sp = sbox.text_frame.paragraphs[0]; sp.text = subtitle; sp.font.size = Pt(22)
            sp.font.name = BODY_FONT; sp.font.color.rgb = RGBColor(255,255,255)
            # author
            abox = slide.shapes.add_textbox(Inches(1.2), Inches(4.1), Inches(11.0), Inches(0.6))
            ap = abox.text_frame.paragraphs[0]; ap.text = author; ap.font.size = Pt(16)
            ap.font.name = BODY_FONT; ap.font.color.rgb = RGBColor(230,230,230)
            # tagline bg for red
            if with_tagline_bg:
                tbg = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.4))
                tp = tbg.text_frame.paragraphs[0]; tp.text = "傳承，不只是資產的延續，而是幸福的流動。"
                tp.font.size = Pt(28); tp.font.name = TITLE_FONT; tp.font.color.rgb = RGBColor(255,255,255)

        elif title == "封底":
            tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.8), Inches(1.2))
            p = tb.text_frame.paragraphs[0]; p.text = "Thank You"
            p.font.size = Pt(40); p.font.bold = True; p.font.name = TITLE_FONT; p.font.color.rgb = RGBColor(255,255,255)

            qbox = slide.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(11.0), Inches(1.6))
            qp = qbox.text_frame.paragraphs[0]; qp.text = closing_quote
            qp.font.size = Pt(20); qp.font.name = BODY_FONT; qp.font.color.rgb = RGBColor(240,240,240)
        else:
            # section
            tbox = slide.shapes.add_textbox(Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.9))
            tp = tbox.text_frame.paragraphs[0]; tp.text = title
            tp.font.size = Pt(30); tp.font.bold = True; tp.font.name = TITLE_FONT
            tp.font.color.rgb = RGBColor(*gold)

            bbox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.2), Inches(2.8))
            btf = bbox.text_frame
            p1 = btf.paragraphs[0]; p1.text = "• " + line1; p1.font.size = Pt(22); p1.font.name = BODY_FONT
            p1.font.color.rgb = RGBColor(255,255,255)
            p2 = btf.add_paragraph(); p2.text = "• " + (line2 or ""); p2.font.size = Pt(22)
            p2.font.name = BODY_FONT; p2.font.color.rgb = RGBColor(230,230,230)

        footer(slide)

    bio = BytesIO(); pres.save(bio); bio.seek(0)
    return bio.read()

def build_pdf_bytes(theme: str, *, author: str, footer_text: str, closing_quote: str) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import Color, white
    from reportlab.lib.units import cm
    import textwrap

    def hex_to_color(hexstr):
        hexstr = hexstr.strip().lstrip('#')
        r,g,b = [int(hexstr[i:i+2],16)/255.0 for i in (0,2,4)]
        return Color(r,g,b)

    bg = hex_to_color(BRAND[theme]["bg"])
    gold = hex_to_color(BRAND[theme]["gold"])

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    W, H = A4

    def footer():
        c.setFillColor(gold); c.rect(0, 0, W, 28, fill=1, stroke=0)
        c.setFont("Helvetica", 9); c.setFillColor(white)
        c.drawString(20, 10, footer_text)

    # cover
    c.setFillColor(bg); c.rect(0,0,W,H,fill=1,stroke=0)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 24); c.drawString(3*cm, H-5*cm, "永傳家族辦公室｜商業策略報告")
    c.setFont("Helvetica", 14); c.drawString(3*cm, H-6*cm, "讓智慧，承載愛的延續。")
    c.setFont("Helvetica", 12); c.drawString(3*cm, H-7*cm, author)
    footer(); c.showPage()

    # content
    c.setFillColor(bg); c.rect(0,0,W,H,fill=1,stroke=0)
    y = H-3*cm
    for title, body in REPORT_PARAGRAPHS:
        c.setFillColor(gold); c.setFont("Helvetica-Bold", 14); c.drawString(3*cm, y, title); y -= 0.9*cm
        c.setFillColor(white); c.setFont("Helvetica", 11)
        for line in textwrap.wrap(body, width=38):
            c.drawString(3*cm, y, line); y -= 0.6*cm
        y -= 0.4*cm
        if y < 4*cm:
            footer(); c.showPage(); c.setFillColor(bg); c.rect(0,0,W,H,fill=1,stroke=0)
            y = H-3*cm

    # outro
    c.setFillColor(bg); c.rect(0,0,W,H,fill=1,stroke=0)
    c.setFillColor(white); c.setFont("Helvetica-Bold", 22); c.drawString(3*cm, H-5*cm, "謝謝！")
    c.setFont("Helvetica", 12)
    for i, line in enumerate(closing_quote.split("\n")):
        c.drawString(3*cm, H-6*cm - i*0.7*cm, line)
    footer(); c.showPage(); c.save()
    buf.seek(0)
    return buf.read()

def build_readme_bytes() -> bytes:
    content = (
        "📘 gracefo_strategy_pack 說明\n\n"
        "本資料包共包含 4 份文件，分為「品牌藍金版」（專業提案風格）與「品牌紅金版」（溫暖敘事風格）。\n\n"
        "檔案內容：\n"
        "1) 永傳家族辦公室｜商業策略藍圖（藍金版 .pptx） → 適合對外簡報、顧問培訓、策略會議\n"
        "2) 永傳家族辦公室｜商業策略藍圖（紅金版 .pptx） → 適合品牌故事分享會、客戶講座\n"
        "3) 永傳家族辦公室｜商業策略報告（藍金版 .pdf） → 投資人與策略夥伴閱讀\n"
        "4) 永傳家族辦公室｜商業策略報告（紅金版 .pdf） → 品牌行銷與內部學習使用\n\n"
        "封底語：\n"
        "「讓智慧，承載愛的延續。\n永傳家族辦公室，與您一起，擘劃幸福的未來。」\n"
    )
    return content.encode("utf-8")

def build_zip(author: str | None = None,
              footer_text: str | None = None,
              closing_quote: str | None = None,
              logo_bytes: bytes | None = None,
              subtitle_blue: str | None = None,
              subtitle_red: str | None = None,
              zip_filename: str = "永傳家族辦公室_商業策略藍圖.zip") -> tuple[str, bytes]:
    """回傳 (zip_filename, zip_bytes)"""
    author = author or DEFAULTS["author"]
    footer_text = footer_text or DEFAULTS["footer_text"]
    closing_quote = closing_quote or DEFAULTS["closing_quote"]
    subtitle_blue = subtitle_blue or DEFAULTS["subtitle_blue"]
    subtitle_red = subtitle_red or DEFAULTS["subtitle_red"]

    # build 4 docs in memory
    blue_ppt = build_pptx_bytes("blue", author=author, footer_text=footer_text,
                                closing_quote=closing_quote, logo_bytes=logo_bytes,
                                subtitle=subtitle_blue, with_tagline_bg=False)
    red_ppt  = build_pptx_bytes("red", author=author, footer_text=footer_text,
                                closing_quote=closing_quote, logo_bytes=logo_bytes,
                                subtitle=subtitle_red, with_tagline_bg=True)
    blue_pdf = build_pdf_bytes("blue", author=author, footer_text=footer_text, closing_quote=closing_quote)
    red_pdf  = build_pdf_bytes("red", author=author, footer_text=footer_text, closing_quote=closing_quote)
    readme   = build_readme_bytes()

    mem = BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("永傳家族辦公室｜商業策略藍圖_藍金版.pptx", blue_ppt)
        z.writestr("永傳家族辦公室｜商業策略藍圖_紅金版.pptx",  red_ppt)
        z.writestr("永傳家族辦公室｜商業策略報告_藍金版.pdf", blue_pdf)
        z.writestr("永傳家族辦公室｜商業策略報告_紅金版.pdf",  red_pdf)
        z.writestr("readme.txt", readme)
    mem.seek(0)
    return zip_filename, mem.read()
